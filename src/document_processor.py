"""
Document Processing Module
Handles various document formats including PDF, DOCX, XLSX, PPTX, etc.
"""
import os
import io
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
import pandas as pd
from PIL import Image
from pdf2image import convert_from_path
import pytesseract

class DocumentProcessor:
    """Process various document formats and extract text."""
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from PDF file."""
        text = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text.append(page.extract_text() or "")
        except Exception as e:
            print(f"Error reading PDF: {e}")
            # Fallback to OCR if text extraction fails
            try:
                images = convert_from_path(file_path)
                for i, image in enumerate(images):
                    text.append(pytesseract.image_to_string(image))
            except Exception as ocr_error:
                print(f"OCR fallback failed: {ocr_error}")
        return "\n".join(text)
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            doc = Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            print(f"Error reading DOCX: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_xlsx(file_path: str) -> str:
        """Extract text from Excel file."""
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True)
            text = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        text.append(" ".join(row_text))
            return "\n".join(text)
        except Exception as e:
            print(f"Error reading XLSX: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from PowerPoint file."""
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            print(f"Error reading PPTX: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_image(file_path: str) -> str:
        """Extract text from image using OCR."""
        try:
            return pytesseract.image_to_string(Image.open(file_path))
        except Exception as e:
            print(f"Error processing image: {e}")
            return ""
    
    @classmethod
    def extract_text_from_file(cls, file_path: str) -> str:
        """Extract text from any supported file type."""
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.pdf':
            return cls.extract_text_from_pdf(file_path)
        elif file_ext == '.docx':
            return cls.extract_text_from_docx(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            return cls.extract_text_from_xlsx(file_path)
        elif file_ext == '.pptx':
            return cls.extract_text_from_pptx(file_path)
        elif file_ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
            return cls.extract_text_from_image(file_path)
        elif file_ext == '.txt':
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except Exception as e:
                print(f"Error reading text file: {e}")
                return ""
        else:
            print(f"Unsupported file format: {file_ext}")
            return ""
